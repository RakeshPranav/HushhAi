import io
import pdfplumber
from pptx import Presentation
import pytesseract
from PIL import Image

def extract_text(content: bytes, filename: str, content_type: str) -> str:
    text = ""
    
    try:
        if content_type == "application/pdf":
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                        
        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(io.BytesIO(content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                        
        elif content_type in ["image/png", "image/jpeg", "image/jpg"]:
            image = Image.open(io.BytesIO(content))
            # Requires tesseract installed on the system
            text = pytesseract.image_to_string(image)
            
        elif content_type == "text/plain":
            text = content.decode("utf-8")
            
        else:
            raise Exception("Unsupported file format for extraction.")
            
    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
        # Return whatever we got or throw
        if not text:
            raise e
            
    # Clean the text a bit
    return text.strip()
